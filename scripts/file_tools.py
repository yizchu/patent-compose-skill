import os
import json
import re
import base64
from pathlib import Path
from pptx import Presentation
from docx import Document
import fitz
import pdfplumber
import win32com.client
import pythoncom
import zipfile

CREDENTIAL_STORE_DIRS = frozenset({
    ".ssh", ".gnupg", ".aws", ".gcloud", "secrets", ".secrets", "credentials",
})
ENV_SUFFIXES = frozenset({".example", ".sample", ".template", ".dist", ".env", ".envrc"})
SENSITIVE_PATTERNS = [
    re.compile(r'(^|[\\/])\.(env|envrc)(\.|$)', re.IGNORECASE),
    re.compile(r'\.(pem|key|p12|pfx|cert|crt|der|p8)$', re.IGNORECASE),
    re.compile(r'(^|[^A-Za-z0-9])(id_rsa|id_dsa|id_ecdsa|id_ed25519)(\.pub)?$', re.IGNORECASE),
    re.compile(r'^secring(\.(gpg|pgp))?$', re.IGNORECASE),
    re.compile(r'(\.netrc|\.pgpass|\.htpasswd|\.npmrc|\.pypirc|\.git-credentials|\.boto)$', re.IGNORECASE),
]
GENERIC_KEYWORD_PATTERNS = [
    re.compile(r'(?<![a-zA-Z0-9])(credential|secret|passwd|password|private_key)s?(?![a-zA-Z])', re.IGNORECASE),
    re.compile(r'(?<![a-zA-Z0-9])tokens?(?![a-zA-Z])', re.IGNORECASE),
    re.compile(r'(?<![a-zA-Z0-9])service[._-]?account(?![a-zA-Z])', re.IGNORECASE),
    # Chinese keywords
    re.compile(r'(?:密码|密钥|口令|凭证|私钥|公钥|令牌|鉴权|认证|授权|证书|敏感信息|访问密钥|安全凭证)'),
]
SECRET_PATTERNS = [
    # API Keys (various providers)
    re.compile(r'(?:api[_-]?key|apikey)\s*[=:]\s*["\']?([A-Za-z0-9]{20,})["\']?', re.IGNORECASE),
    re.compile(r'(?:sk|pk)[_-](?:live|test|prod|dev)[_-][A-Za-z0-9]{20,}', re.IGNORECASE),

    # AWS Credentials
    re.compile(r'AKIA[0-9A-Z]{16}', re.IGNORECASE),  # AWS Access Key ID
    re.compile(r'(?:aws[_-]?)?(?:secret[_-]?access[_-]?key|aws[_-]?secret[_-]?key)\s*[=:]\s*["\']?([A-Za-z0-9/+=]{40})["\']?', re.IGNORECASE),

    # Azure Credentials
    re.compile(r'(?:azure|aad)[_-]?(?:client[_-]?secret|tenant[_-]?id)\s*[=:]\s*["\']?[^\s"\'<>]{10,}["\']?', re.IGNORECASE),

    # Google Cloud Credentials
    re.compile(r'"project_id"\s*:\s*"[^"]+"', re.IGNORECASE),
    re.compile(r'"private_key_id"\s*:\s*"[^"]+"', re.IGNORECASE),

    # Alibaba Cloud AccessKey
    re.compile(r'LTAI[A-Za-z0-9]{12,20}'),

    # Tencent Cloud
    re.compile(r'(?:secret[_-]?id|secret[_-]?key)\s*[=:]\s*["\']?[A-Za-z0-9]{20,}["\']?', re.IGNORECASE),

    # Generic API tokens
    re.compile(r'(?:bearer|token|auth)\s+(?:[A-Za-z0-9\-_]{20,})', re.IGNORECASE),
    re.compile(r'(?:x[_-]?api[_-]?key|api[_-]?token)\s*[=:]\s*["\']?([A-Za-z0-9]{16,})["\']?', re.IGNORECASE),

    # Database connection strings
    re.compile(r'(?:mongodb|postgres|mysql|redis|amqp)://[^\s"\'<>]{10,}', re.IGNORECASE),
    re.compile(r'(?:server|host|database)\s*[=:]\s*["\']?[^\s"\'<>]{5,}', re.IGNORECASE),

    # Passwords
    re.compile(r'(?:password|passwd|pwd)\s*[=:]\s*["\']?[^\s"\'<>]{8,}["\']?', re.IGNORECASE),

    # Private keys (PEM format)
    re.compile(r'-----BEGIN (?:RSA |EC |DSA )?PRIVATE KEY-----'),
    re.compile(r'-----BEGIN OPENSSH PRIVATE KEY-----'),
    re.compile(r'-----BEGIN PGP PRIVATE KEY BLOCK-----'),

    # OAuth/Service Account
    re.compile(r'"type"\s*:\s*"service_account"', re.IGNORECASE),
    re.compile(r'(?:oauth|client)[_-]?(?:secret|token)\s*[=:]\s*["\']?[A-Za-z0-9\-_]{16,}["\']?', re.IGNORECASE),

    # GitHub/GitLab tokens
    re.compile(r'ghp_[A-Za-z0-9]{36}'),  # GitHub Personal Access Token
    re.compile(r'glpat-[A-Za-z0-9\-]{20,}'),  # GitLab Personal Access Token
    re.compile(r'github_pat_[A-Za-z0-9]{22}_[A-Za-z0-9]{59}'),  # GitHub Fine-grained PAT

    # Slack tokens
    re.compile(r'xox[baprs]-[0-9]{10,13}-[a-zA-Z0-9-]+'),

    # Stripe payment keys
    re.compile(r'sk_(?:live|test)_[A-Za-z0-9]{24,}'),
    re.compile(r'rk_(?:live|test)_[A-Za-z0-9]{24,}'),

    # Twilio credentials
    re.compile(r'AC[a-f0-9]{32}'),  # Account SID
    re.compile(r'SK[a-f0-9]{32}'),  # API Key

    # SendGrid API key
    re.compile(r'SG\.[A-Za-z0-9_-]{22,}\.[A-Za-z0-9_-]{43,}'),

    # Firebase API Key
    re.compile(r'AIza[0-9A-Za-z_-]{35}'),

    # JWT tokens
    re.compile(r'eyJ[A-Za-z0-9_-]{10,}\.[A-Za-z0-9_-]{10,}\.[A-Za-z0-9_-]{10,}'),

    # Base64 encoded credentials
    re.compile(r'(?:key|secret|token|password|credential)\s*[=:]\s*["\']?([A-Za-z0-9+/]{40,}={0,2})["\']?', re.IGNORECASE),

    # Webhook URLs with tokens
    re.compile(r'https?://[^\s]*(?:hook|callback|webhook)[^\s]*[?&](?:token|key|secret)=[^\s&]{10,}', re.IGNORECASE),

    # CI/CD credentials
    re.compile(r'(?:docker[_-]?password|registry[_-]?token)\s*[=:]\s*["\']?[^\s"\'<>]{8,}["\']?', re.IGNORECASE),
    re.compile(r'(?:k8s|kube)[_-]?(?:token|secret)\s*[=:]\s*["\']?[A-Za-z0-9._-]{20,}["\']?', re.IGNORECASE),
    re.compile(r'jenkins[_-]?(?:token|password|secret)\s*[=:]\s*["\']?[^\s"\'<>]{10,}["\']?', re.IGNORECASE),

    # Session IDs and cookies
    re.compile(r'(?:session[_-]?id|sid|phpsessid|jsessionid)\s*[=:]\s*["\']?[A-Za-z0-9]{16,}["\']?', re.IGNORECASE),
    re.compile(r'(?:auth[_-]?token|access[_-]?token|refresh[_-]?token)=([A-Za-z0-9._-]{20,})', re.IGNORECASE),

    # Mobile signing keys
    re.compile(r'(?:signing[_-]?key|store[_-]?password|key[_-]?password)\s*[=:]\s*["\']?[^\s"\'<>]{6,}["\']?', re.IGNORECASE),

    # SSH host keys and fingerprints
    re.compile(r'(?:ecdsa-sha2-nistp256|ssh-rsa|ssh-ed25519)\s+[A-Za-z0-9+/=]{20,}'),

    # Windows product keys
    re.compile(r'[A-Z0-9]{5}-[A-Z0-9]{5}-[A-Z0-9]{5}-[A-Z0-9]{5}-[A-Z0-9]{5}'),

    # Generic high-entropy assignments (catches custom secret formats)
    re.compile(r'(?:secret|key|token|credential)\s*[=:]\s*["\']([A-Za-z0-9+/=_-]{32,})["\']', re.IGNORECASE),

    # Chinese secret patterns (中文密钥模式)
    re.compile(r'(?:密码|口令|pwd)\s*[=:：]\s*["\']?[^\s"\'<>]{8,}["\']?'),
    re.compile(r'(?:密钥|私钥|公钥|access[_-]?key|secret[_-]?key)\s*[=:：]\s*["\']?[A-Za-z0-9+/=_-]{16,}["\']?'),
    re.compile(r'(?:令牌|token|鉴权码)\s*[=:：]\s*["\']?[A-Za-z0-9._-]{16,}["\']?'),
    re.compile(r'(?:凭证|credential|证书)\s*[=:：]\s*["\']?[^\s"\'<>]{10,}["\']?'),
    re.compile(r'(?:数据库连接|数据库地址|db[_-]?url|conn[_-]?str)\s*[=:：]\s*["\']?[^\s"\'<>]{10,}["\']?'),
    re.compile(r'(?:服务器|主机|host|地址)\s*[=:：]\s*["\']?[^\s"\'<>]{5,}["\']?'),
]
PII_PATTERNS = [
    # Email addresses
    re.compile(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'),

    # Phone numbers (various formats)
    re.compile(r'(?:\+?86)?1[3-9]\d{9}'),  # Chinese phone numbers
    re.compile(r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b'),  # US phone numbers
    re.compile(r'\+\d{1,3}[-.\s]?\(?\d{1,4}\)?[-.\s]?\d{1,4}[-.\s]?\d{1,9}'),  # International

    # Chinese ID card numbers (18 digits)
    re.compile(r'\b\d{17}[\dXx]\b'),

    # Social Security Numbers (US)
    re.compile(r'\b\d{3}-\d{2}-\d{4}\b'),

    # Credit card numbers (basic pattern, not validated)
    re.compile(r'\b(?:\d{4}[-\s]?){3}\d{4}\b'),

    # IP addresses (private/internal)
    re.compile(r'\b(?:10|172\.(?:1[6-9]|2\d|3[01])|192\.168)\.\d{1,3}\.\d{1,3}\b'),

    # MAC addresses
    re.compile(r'(?:[0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}'),

    # Internal domain names
    re.compile(r'(?:[\w-]+\.)+(?:internal|local|corp|intranet|private|lan)(?:\.[a-z]{2,})?', re.IGNORECASE),

    # VPN configurations
    re.compile(r'(?:vpn|openvpn|wireguard)[_-]?(?:key|secret|config)\s*[=:]\s*["\']?[^\s"\'<>]{10,}["\']?', re.IGNORECASE),

    # Proxy server credentials
    re.compile(r'(?:proxy|http_proxy|https_proxy)\s*[=:]\s*["\']?http[s]?://[^\s"\'<>:]+:[^\s"\'<>]+@[^\s"\'<>]+["\']?', re.IGNORECASE),

    # Windows credentials
    re.compile(r'(?:windows[_-]?credential|ntlm[_-]?hash|kerberos[_-]?key)\s*[=:]\s*["\']?[^\s"\'<>]{10,}["\']?', re.IGNORECASE),

    # Environment variable references with sensitive names
    re.compile(r'\$\{?(?:SECRET|PASSWORD|TOKEN|API_KEY|PRIVATE_KEY)[_A-Z]*\}?', re.IGNORECASE),

    # Chinese PII patterns (中文个人信息)
    re.compile(r'(?:姓名|名字|真实姓名|联系人)\s*[=:：]\s*[\u4e00-\u9fa5]{2,4}'),
    re.compile(r'(?:地址|住址|家庭住址|详细地址|通讯地址)\s*[=:：]\s*[\u4e00-\u9fa50-9a-zA-Z#\-]{10,}'),
    re.compile(r'(?:身份证|身份证号|身份证号码|证件号)\s*[=:：]\s*\d{17}[\dXx]'),
    re.compile(r'(?:邮箱|电子邮箱|邮件地址)\s*[=:：]\s*[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'),
    re.compile(r'(?:手机|手机号|联系电话|联系方式|电话号码)\s*[=:：]\s*(?:\+?86)?1[3-9]\d{9}'),
    re.compile(r'(?:银行卡|银行卡号|银行卡号码)\s*[=:：]\s*\d{16,19}'),
    re.compile(r'(?:护照|护照号|护照号码)\s*[=:：]\s*[A-Za-z0-9]{8,12}'),
]
SENSITIVE_KEY_PATTERNS = [
    # 匹配 key = value, key: value, key value 等格式
    re.compile(r'((?:api[_-]?key|apikey|secret[_-]?key|access[_-]?key|private[_-]?key|password|passwd|pwd|token|auth[_-]?token|bearer|credential)[\s:=]+)["\']?([^\s"\']{8,})["\']?', re.IGNORECASE),
    # 匹配数据库连接信息
    re.compile(r'((?:connection[_-]?string|conn[_-]?str|database[_-]?url|db[_-]?url)[\s:=]+)["\']?([^\s"\']{10,})["\']?', re.IGNORECASE),
    # 匹配 AWS secret access key
    re.compile(r'((?:aws[\s_-]*)?(?:secret[\s_-]+access[\s_-]+key)[\s:=]+)["\']?([A-Za-z0-9/+=]{20,})["\']?', re.IGNORECASE),
    # 匹配中文密钥格式 (密码 = xxx, 密钥：xxx, 令牌: xxx)
    re.compile(r'((?:密码|密钥|口令|凭证|令牌|鉴权码|私钥|公钥|证书|数据库密码|数据库密钥|访问密钥)[\s:=：]+)["\']?([^\s"\']{8,})["\']?'),
    # 匹配无分隔符的中文密钥格式 (密码xxx)，排除后面跟中文字符的情况（如"密码学"、"密码保护"）
    re.compile(r'((?:密码|密钥|口令|凭证|令牌|鉴权码|私钥|公钥|证书|数据库密码|数据库密钥|访问密钥))(?![\u4e00-\u9fa5\s:=：])["\']?([A-Za-z0-9+/=_-]{8,})["\']?'),
]

def to_json(data: dict, file_path: str):
    with open(file_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

def from_json(file_path: str):
    with open(file_path, 'r', encoding='utf-8') as f:
        return json.load(f)

def clean_filename(filename: str) -> str:
    """清理文件名中的无效字符，返回安全的文件名。 \n
    Windows 文件名无效字符: <>:"/\|?* 以及控制字符
    """
    invalid_chars = set('<>:"/\\|?*')
    # 将无效字符替换为下划线
    cleaned = ''.join(c if c not in invalid_chars else '_' for c in filename)
    # 移除控制字符（ASCII 0-31）
    cleaned = ''.join(c for c in cleaned if ord(c) >= 32)
    # 移除首尾空格和点
    cleaned = cleaned.strip('. ')
    # 如果文件名为空，返回默认名称
    if not cleaned:
        cleaned = 'unnamed_file'
    # 限制文件名长度（Windows 最大255字符）
    if len(cleaned) > 200:
        cleaned = cleaned[:200]
    return cleaned

def is_sensitive_file(file_path: Path|str) -> bool:
    if isinstance(file_path, str):
        file_path = Path(file_path)
    parent = file_path.parts[:-1]
    if any(part.lower() in CREDENTIAL_STORE_DIRS for part in parent):
        return True
    name = file_path.name
    if any(p.search(name) for p in SENSITIVE_PATTERNS):
        return True
    if any(p in name.lower() for p in ENV_SUFFIXES):
        return True

    def generic_keyword_hit(name: str) -> bool:
        """True if a generic secret keyword appears load-bearing in the filename.

        Secret-store files name their contents, and in English compounds the
        content noun is the head, which comes last: "github-personal-access-token",
        "api_token", "oauth_token". A keyword that is neither at the end of the
        stem nor in a short (<=2 word) name is a topic word in a descriptive slug
        ("token-economics-of-recall.md", "password-policy-discussion.md") and must
        not cause the file to be silently dropped from the graph (#436, #718).
        """
        # Stem = name minus only the FINAL extension (not up to the first dot), so a
        # multi-dot topic slug like `token.economics.notes.md` keeps all its words and
        # doesn't collapse to a bare `token` (#2106). Leading dots stripped so
        # dotfiles like `.token` keep their keyword.
        stem = Path(name).stem.lstrip('.') or Path(name).stem
        for pat in GENERIC_KEYWORD_PATTERNS:
            hit = False
            for m in pat.finditer(stem):
                hit = True
                if m.end() == len(stem):  # keyword ends the stem -> names the contents
                    return True
            if hit and len([w for w in re.compile(r'[-_\s.]+').split(stem) if w]) <= 2:
                return True  # short name like token_config.yaml / secret_handler.txt
        return False

    return generic_keyword_hit(name)

def content_remove_sensitive(content: str) -> str:
    # 优先移除完整的 PEM 私钥块（必须在 SECRET_PATTERNS 之前，因为 SECRET_PATTERNS 会单独匹配 BEGIN 行，破坏块结构）
    content = re.sub(
        r'-----BEGIN [A-Z ]*PRIVATE KEY-----.*?-----END [A-Z ]*PRIVATE KEY-----',
        '[REDACTED]',
        content,
        flags=re.DOTALL
    )
    # 去除密钥/凭证类敏感信息
    for p in SECRET_PATTERNS:
        content = p.sub('[REDACTED]', content)
    # 去除个人信息（PII）
    for p in PII_PATTERNS:
        content = p.sub('[REDACTED]', content)
    # 清理可能残留的敏感键值对（处理 Word/PPT 等非结构化文本）
    for p in SENSITIVE_KEY_PATTERNS:
        content = p.sub(lambda m: m.group(1) + '[REDACTED]', content)
    return content


# ============================================================
# 文件文本和图片提取器
# ============================================================
class FileReader:
    def __init__(self, file_path: str, extension: str, out_dir: str):
        self.file_path = file_path
        self.extension = extension.replace(".", "")
        self.out_dir = out_dir
        os.makedirs(self.out_dir, exist_ok=True)

    def read_texts(self) -> str:
        """读取文件的文本内容，去除敏感信息"""
        text_output = None
        if self.extension in {'pptx', 'ppt', 'pptm'}:
            text_output = self.get_ppt_text()
        elif self.extension in {'docx', 'doc', 'docm'} and self._zip_within_caps():
            text_output = self.get_word_text()
        elif self.extension in {'pdf'}:
            text_output = self.get_pdf_text()
        elif self.extension in {'ipynb'}:
            text_output = self.get_jupyter_text()
        if text_output is not None:
            text_output = content_remove_sensitive(text_output)
            with open(os.path.join(self.out_dir, "content.md"), 'w', encoding='utf-8', errors='ignore') as f:
                f.write(text_output)
            return text_output
        else:
            for encoding in ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'latin-1']:
                try:
                    with open(self.file_path, 'r', encoding=encoding, errors='ignore') as f_input:
                        text_output = content_remove_sensitive(f_input.read())
                    with open(os.path.join(self.out_dir, "content.md"), 'w', encoding=encoding, errors='ignore') as f_output:
                        f_output.write(text_output)
                    return text_output
                except:
                    pass
            return ""

    def read_images(self) -> None:
        """提取文件中的图片"""
        if self.extension in {'pptx', 'ppt', 'pptm'}:
            return self.get_ppt_image()
        elif self.extension in {'docx', 'doc', 'docm'} and self._zip_within_caps():
            return self.get_word_image()
        elif self.extension in {'pdf'}:
            return self.get_pdf_image()
        elif self.extension in {'ipynb'}:
            return self.get_jupyter_image()

    def _zip_within_caps(self) -> bool:
        """过滤掉过大的基于 zip 的 Office 文件。
        """
        path = Path(self.file_path)

        try:
            if path.stat().st_size > 50 * 1024 * 1024:
                return False
        except OSError:
            return False

        try:
            with zipfile.ZipFile(path) as zf:
                infos = zf.infolist()
                compressed = sum(i.compress_size for i in infos) or 1
                declared = sum(i.file_size for i in infos)
                if declared > 512 * 1024 * 1024:
                    return False
                if declared / compressed > 200:
                    return False
                total = 0
                for info in infos:
                    with zf.open(info) as member:
                        while True:
                            chunk = member.read(1024 * 1024)
                            if not chunk:
                                break
                            total += len(chunk)
                            if total > 512 * 1024 * 1024:
                                return False
        except (zipfile.BadZipFile, OSError, EOFError):
            return False
        return True

    def get_ppt_text(self) -> str:
        text_output = []
        if self.extension == "ppt":
            try:
                pythoncom.CoInitialize()
                app = win32com.client.Dispatch("PowerPoint.Application")
                prs = app.Presentations.Open(self.file_path, WithWindow=False)
                for slide_number, slide in enumerate(prs.Slides, start=1):
                    slide_text = []
                    for shape in slide.Shapes:
                        if shape.HasTextFrame and shape.TextFrame.HasText:
                            text = shape.TextFrame.TextRange.Text.strip()
                            if text:
                                slide_text.append(text)
                    if slide_text:
                        text_output.append(f"### Slide {slide_number}\n" + "\n".join(slide_text).strip())
                prs.Close()
                app.Quit()
                pythoncom.CoUninitialize()
            except:
                if 'app' in locals():
                    app.Quit()
                pythoncom.CoUninitialize()
                return ""
        else:
            try:
                presentation = Presentation(self.file_path)
                for slide_number, slide in enumerate(presentation.slides, start=1):
                    slide_text = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                slide_text.append("".join(run.text for run in paragraph.runs))
                    if slide_text:
                        text_output.append(f"### Slide {slide_number}\n" + "\n".join(slide_text).strip())
            except:
                return ""
        text_output = "\n".join(text_output)
        return text_output

    def get_ppt_image(self) -> None:
        image_cnt = 1
        if self.extension == "ppt":
            try:
                pythoncom.CoInitialize()
                app = win32com.client.Dispatch("PowerPoint.Application")
                prs = app.Presentations.Open(self.file_path, WithWindow=False)
                for slide in prs.Slides:
                    for shape in slide.Shapes:
                        if shape.Type == 13:
                            try:
                                shape.Export(os.path.join(self.out_dir, f"image{image_cnt}.png"), 2)
                                image_cnt += 1
                            except:
                                pass
                prs.Close()
                app.Quit()
                pythoncom.CoUninitialize()
            except Exception as e:
                if 'app' in locals():
                    app.Quit()
                pythoncom.CoUninitialize()
        else:
            presentation = Presentation(self.file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    try:
                        if "image" in shape.image.content_type:
                            with open(os.path.join(self.out_dir, f"image{image_cnt}.{shape.image.ext}"), 'wb') as f:
                                f.write(shape.image.blob)
                            image_cnt += 1
                    except:
                        pass

    def get_word_text(self) -> str:
        '''仅支持 .docx 格式'''
        paragraphs = []
        try:
            doc = Document(self.file_path)
            for para in doc.paragraphs:
                # 跳过表格
                if para._element.getparent().tag.endswith('tbl') or para._element.getparent().tag.endswith('tc'):
                    continue
                if para.text.strip():
                    style = para.style.name if para.style else ""
                    if "Heading" in style or "标题" in style:
                        level = style.replace("Heading ", "").replace("标题 ", "")
                        prefix = "#" * int(level) if level.isdigit() else "##"
                        paragraphs.append(f"{prefix} {para.text}")
                    else:
                        paragraphs.append(para.text)
        except:
            return ""
        return '\n'.join(paragraphs)

    def get_word_image(self) -> None:
        if self.extension == "doc":
            try:
                pythoncom.CoInitialize()
                app = win32com.client.DispatchEx("Word.Application")
                app.Visible = False
                doc = app.Documents.Open(self.file_path)
                # 将.doc另存为临时.docx，然后提取图片
                temp_docx = os.path.join(self.out_dir, "_temp_extract.docx")
                doc.SaveAs2(temp_docx, FileFormat=16)  # 16 = wdFormatXMLDocument
                # 从临时docx中提取图片
                with zipfile.ZipFile(temp_docx, 'r') as z:
                    for item in z.namelist():
                        try:
                            if item.startswith('word/media/'):
                                with open(os.path.join(self.out_dir, os.path.basename(item)), 'wb') as f:
                                    f.write(z.read(item))
                        except:
                            pass
                # 清理临时文件
                if os.path.exists(temp_docx):
                    os.remove(temp_docx)
                doc.Close()
                app.Quit()
                pythoncom.CoUninitialize()
            except Exception as e:
                if 'app' in locals():
                    app.Quit()
                pythoncom.CoUninitialize()
        else:
            with zipfile.ZipFile(self.file_path, 'r') as z:
                for item in z.namelist():
                    try:
                        if item.startswith('word/media/'):
                            with open(os.path.join(self.out_dir, os.path.basename(item)), 'wb') as f:
                                f.write(z.read(item))
                    except:
                        pass

    def get_pdf_text(self) -> str:
        pages_text = []
        try:
            with pdfplumber.open(self.file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        pages_text.append(f"## Page {i}\n{text}")
            return '\n'.join(pages_text)
        except:
            return ""

    def get_pdf_image(self) -> None:
        image_cnt = 1
        doc = fitz.open(self.file_path)
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            image_list = page.get_images(full=True)
            for image in image_list:
                try:
                    xref = image[0]
                    base_image = doc.extract_image(xref)
                    image_path = os.path.join(self.out_dir, f"image{image_cnt}.{base_image['ext']}")
                    image_cnt += 1
                    with open(image_path, "wb") as f:
                        f.write(base_image["image"])
                except:
                    pass
        doc.close()

    def get_jupyter_text(self) -> str:
        cells_text = []
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                notebook = json.load(f)
            for i, cell in enumerate(notebook.get('cells', [])):
                cell_type = cell.get('cell_type', 'unknown')
                source = ''.join(cell.get('source', []))
                if source.strip():
                    if cell_type == 'code':
                        cells_text.append(f"## Cell {i+1} (Code)\n```python\n{source}\n```")
                    elif cell_type == 'markdown':
                        cells_text.append(f"## Cell {i+1} (Markdown)\n{source}")
                    elif cell_type == 'raw':
                        cells_text.append(f"## Cell {i+1} (Raw)\n{source}")
                    else:
                        cells_text.append(f"## Cell {i+1}\n{source}")
            return '\n'.join(cells_text)
        except:
            return ""

    def get_jupyter_image(self) -> None:
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                    notebook = json.load(f)
            for i, cell in enumerate(notebook.get('cells', [])):
                if cell.get('cell_type', 'unknown') == 'code':
                    image_cnt = 1
                    outputs = cell.get('outputs', [])
                    for output in outputs:
                        if 'data' in output:
                            data = output['data']
                            if 'image/png' in data:
                                png_data = base64.b64decode(data['image/png'])
                                image_path = os.path.join(self.out_dir, f"cell{i+1}_image{image_cnt}.png")
                                image_cnt += 1
                                with open(image_path, 'wb') as f:
                                    f.write(png_data)
                            if 'image/jpeg' in data:
                                jpeg_data = base64.b64decode(data['image/jpeg'])
                                image_path = os.path.join(self.out_dir, f"cell{i+1}_image{image_cnt}.jpg")
                                image_cnt += 1
                                with open(image_path, 'wb') as f:
                                    f.write(jpeg_data)
        except:
            pass